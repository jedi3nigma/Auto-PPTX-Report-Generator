from pptx import Presentation
from pptx.util import Inches, Pt, Emu


def gen_test_layout_file(master):
    prs = Presentation(master)
    for i in range(len(prs.slide_layouts)):
        layout = prs.slide_layouts[i]
        slide = prs.slides.add_slide(layout)
        try:
            title = slide.shapes.title
            title.text = 'Title - layout {}'.format(i)
        except:
            pass

        try:
            subtitle = slide.placeholders[1]
            subtitle.text = 'Subtitle - layout {}'.format(i)
        except:
            pass

        for j in range(len(slide.placeholders)):
            try:
                tmp = slide.placeholders[j]
                tmp.text = 'Placeholder {}: {}: layout {}'.format(
                    j, tmp.name, i)
                print(tmp, j, tmp.name, i)
            except:
                continue
    prs.save('slide_layouts_file.pptx')


class SlideSelect():
    _left = Inches(1.5)
    _top = Inches(2)
    _adj_left = Inches(1.25)
    _adj_top = Inches(0.25)

    def __init__(self, out_file, templ_path=None):
        self.out_file = out_file

        if templ_path != None:
            self.pptx = Presentation(templ_path)
        else:
            self.pptx = Presentation()

    def save_file(self):
        self.pptx.save(self.out_file)

    def _set_layout(self, slide_layout):
        if slide_layout == 'title':
            return self.pptx.slide_layouts[0]
        elif slide_layout == 'blank':
            return self.pptx.slide_layouts[1]
        elif slide_layout == 'text':
            return self.pptx.slide_layouts[1]
        elif slide_layout == 'two_columns':
            return self.pptx.slide_layouts[4]
        else:
            return

    def create_slide(self, slide_type, slide_layout, **kwargs):
        layout = self._set_layout(slide_layout)
        if slide_type == 'title':
            return TitleSlide(self.pptx)
        elif slide_type == 'summary':
            return SummarySlide(self.pptx, layout)
        elif slide_type == 'charts':
            return ChartSlide(self.pptx, layout)
        elif slide_type == 'datatable':
            return DataSlide(self.pptx, layout)

    @ staticmethod
    def _set_title(slide, label, fontsize):
        title = slide.shapes.title
        title.text = label
        title.text_frame.paragraphs[0].font.size = Pt(fontsize)
        return title

    @ staticmethod
    def _set_subtitle(slide, label, placehold_num, fontsize):
        subtitle = slide.placeholders[placehold_num]
        subtitle.text = label
        subtitle.text_frame.paragraphs[0].font.size = Pt(fontsize)
        return subtitle

    def _add_image(self, slide, img_path, left, top):
        img = slide.shapes.add_picture(img_path, left, top)
        return img


class TitleSlide(SlideSelect):
    def __init__(self, ppt_obj):
        self.slide = ppt_obj.slides[0]

    def _set_title(self, label, fontsize=30):
        self.title = super()._set_title(self.slide, label, fontsize)

    def _set_subtitle(self, label, placehold_num, fontsize=20):
        self.subtitle = super()._set_subtitle(
            self.slide, label, placehold_num, fontsize)

    def _add_logo(self, path, dimensions):
        left = Inches(dimensions[0])
        top = Inches(dimensions[1])
        img = super()._add_image(self.slide, path, left, top)

    def create(self, **kwargs):
        label = kwargs.get('label', '')
        sub_label = kwargs.get('sub_label', '')
        placehold_num = kwargs.get('placehold_num', 1)
        logo_path_1 = kwargs.get('logo_path_1', '')
        logo_path_2 = kwargs.get('logo_path_2', '')
        logo_dim_1 = kwargs.get('logo_dim_1', [10, 2])
        logo_dim_2 = kwargs.get('logo_dim_2', [8, 2])
        self._set_title(label)
        self._set_subtitle(sub_label, placehold_num)
        self._add_logo(logo_path_1, logo_dim_1)
        self._add_logo(logo_path_2, logo_dim_2)


class SummarySlide(SlideSelect):
    def __init__(self, ppt_obj, add_layout):
        self.slide = ppt_obj.slides.add_slide(add_layout)

    def _set_title(self, label, fontsize=30):
        super()._set_title(self.slide, label, fontsize)

    def _set_summary(self, text, placehold_num):
        summary = self.slide.placeholders[placehold_num]
        summary.text = text
        summary.text_frame.paragraphs[0].font.size = Pt(14)

    def create(self, **kwargs):
        label = kwargs.get('label', '')
        text = kwargs.get('text', '')
        placehold_num = kwargs.get('placehold_num', 1)
        self._set_title(label)
        self._set_summary(text, placehold_num)


class ChartSlide(SlideSelect):
    def __init__(self, ppt_obj, add_layout):
        self.slide = ppt_obj.slides.add_slide(add_layout)

    def _set_title(self, label, fontsize=30):
        super()._set_title(self.slide, label, fontsize)

    def _add_chart(self, path, dimensions, **kwargs):
        if 'boxplt' in kwargs.keys():
            left = dimensions[0] + Inches(kwargs['boxplot'][0])
            top = dimensions[1] + Inches(kwargs['boxplot'][1])
        else:
            left, top = dimensions[0], dimensions[1]
        return super()._add_image(self.slide, path,
                                  left, top)

    def create(self, path, **kwargs):
        label = kwargs.get('label', '')
        dim = kwargs.get('dim', [self._left, self._top])
        self._set_title(label)
        img = self._add_chart(path, dim, **kwargs)


class DataSlide(SlideSelect):
    def __init__(self, ppt_obj, add_layout):
        self.slide = ppt_obj.slides.add_slide(add_layout)

    def _set_title(self, label, fontsize=30):
        super()._set_title(self.slide, label, fontsize)

    def _set_subtitle(self, label, placehold_num, fontsize=20):
        self.subtitle = super()._set_subtitle(
            self.slide, label, placehold_num, fontsize)

    def _add_table_img(self, path, dimensions, append_dim, *args):
        if append_dim:
            left = dimensions[0] + Inches(args[0])
            top = dimensions[1] + Inches(args[1])
        else:
            left, top = dimensions[0], dimensions[1]
        return super()._add_image(self.slide, path,
                                  left, top)

    def _change_table_width(self):
        pass

    def create(self, append_dim=True, **kwargs):
        label = kwargs.get('label', '')
        sub_label = kwargs.get('sub_label', '')
        path = kwargs.get('path', '')
        placehold_num = kwargs.get('placehold_num', 1)
        dim = kwargs.get('dim', [self._left, self._top])
        self._set_title(label)

        for key, table_prop in kwargs.get('table_props').items():
            if key == 'subtitle_prop':
                self._set_subtitle(table_prop[0], table_prop[1])

            if key == 'table_prop':
                img = self._add_table_img(
                    table_prop[0], dim, append_dim, *table_prop[1])
                if label == 'Top Product and Class Revenue':
                    img.width = img.width - Emu(1000000)
